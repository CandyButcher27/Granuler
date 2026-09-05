"""Per-image replacement brief for the template's stock photography.

The template still carries the photography from the deck it was cloned from.
Until a replacement set exists, this hands the assessor one page per image: the
slide it sits on, the exact pixel size to render at, and a prompt to paste into
whichever image generator he prefers. He drops the result into the slide by
hand.

Only substantial photographs are listed. The template's other ~99 images are
icons, rules and SVG glyphs that carry no client identity and stay as they are.
"""
from pathlib import Path

from pptx import Presentation
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, Spacer

from .pdf_generator import STYLES, _build, _heading, _table

TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "assets" / "granuler_template.pptx"

# A picture is a photograph worth replacing, rather than an icon or a rule, when
# it is both physically large on the slide and heavy on disk.
MIN_AREA_SQ_IN = 4.0
MIN_BYTES = 300_000

RENDER_DPI = 150

STYLE_DIRECTION = (
    "Clean modern corporate photography, bright and airy, plenty of negative space, "
    "shallow depth of field. Cool teal and soft green accents on a light neutral ground, "
    "matching a teal-to-green brand gradient. No text, no logos, no watermarks, "
    "no recognisable faces or company branding."
)


def _slide_title(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().splitlines()[0][:80]
    return ""


def collect_photos(template_path: Path | None = None) -> list[dict]:
    return photos_in(Presentation(str(template_path or TEMPLATE_PATH)))


def photos_in(presentation) -> list[dict]:
    """The replaceable photographs in an open deck.

    Takes a Presentation rather than a path so it can run against the deck a
    client actually receives. Slide numbers there differ from the template's -
    gated slides are removed - and five of the titles are rewritten per client,
    so neither is a reliable way to point at a slide from outside the file.
    """
    photos = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.__class__.__name__ != "Picture":
                continue
            try:
                blob = shape.image.blob
            except ValueError:
                continue
            width_in = shape.width / 914400
            height_in = shape.height / 914400
            if width_in * height_in < MIN_AREA_SQ_IN or len(blob) < MIN_BYTES:
                continue
            photos.append({
                "number": len(photos) + 1,
                "slide": slide_number,
                "title": _slide_title(slide),
                "width_in": round(width_in, 2),
                "height_in": round(height_in, 2),
                "width_px": round(width_in * RENDER_DPI),
                "height_px": round(height_in * RENDER_DPI),
            })
    return photos


def _orientation(photo: dict) -> str:
    ratio = photo["width_in"] / photo["height_in"]
    if ratio > 1.15:
        return "wide landscape"
    if ratio < 0.87:
        return "tall portrait"
    return "square"


# Some template pictures are diagrams carrying data, not stock photography, so
# a "photograph of ..." prompt is wrong for them. Keyed by template slide.
SPECIAL_PROMPTS = {
    6: (
        "A clean isometric staircase diagram of five technology maturity bands, dark "
        "background, thin light outlines, rising left to right: \"0-25 Critical Risk\", "
        "\"26-40 Fragile Environment\", \"41-60 Developing Environment\", "
        "\"61-75 Managed Environment\", \"76-90 Mature Digital Operations\", "
        "\"91-100 Digital Leader\". Highlight ONLY the {band_range} step in bright green "
        "and label that step \"{company}\". Every other step stays dim and unlabelled. "
        "No other company name anywhere in the image."
    ),
}

# The steps drawn on the diagram, which are its own bands rather than the
# scoring bands in pptx_generator - the artwork splits the bottom of the range
# in two where the score calculation does not.
_BANDS = [(26, "0-25"), (41, "26-40"), (61, "41-60"), (76, "61-75"), (91, "76-90")]


def band_range_for(score: float | None) -> str:
    if score is None:
        return "41-60"
    for edge, label in _BANDS:
        if score < edge:
            return label
    return "91-100"


def build_prompt(photo: dict, company: str = "", score: float | None = None) -> str:
    special = SPECIAL_PROMPTS.get(photo.get("slide"))
    if special:
        return special.format(
            company=company or "the client company",
            band_range=band_range_for(score),
        )
    subject = photo["title"] or "strategic technology advisory"
    about = f" for {company}" if company else ""
    return (
        f"A photograph illustrating \"{subject}\"{about}, for a technology maturity "
        f"assessment presentation. {STYLE_DIRECTION} "
        f"Compose for a {_orientation(photo)} crop."
    )


def image_brief_pdf(photos: list[dict] | None = None, company: str = "",
                    score: float | None = None) -> bytes:
    photos = collect_photos() if photos is None else photos
    story = _heading(
        "Image Replacement Brief",
        company or "Granuler assessment template",
        "One prompt per stock photograph still carried by the template",
    )

    if not photos:
        story.append(Paragraph("No replaceable photographs were found in the template.",
                               STYLES["body"]))
        return _build(story, "Image Replacement Brief", "Image Replacement Brief · Granuler")

    story.append(Paragraph(
        f"{len(photos)} photographs need replacing. Generate each at the pixel size given, "
        "then drop it into the numbered slide in place of the existing picture. Every other "
        "image in the template is an icon or rule and needs no change. Slide numbers are the "
        "template's, not the generated deck's, which is shorter because unused slides are removed.",
        STYLES["body"],
    ))
    story.append(Spacer(1, 10))

    rows = [[
        Paragraph("#", STYLES["cellhead"]),
        Paragraph("SLIDE", STYLES["cellhead"]),
        Paragraph("SIZE", STYLES["cellhead"]),
        Paragraph("PROMPT", STYLES["cellhead"]),
    ]]
    for photo in photos:
        rows.append([
            Paragraph(f"Image {photo['number']}", STYLES["cell"]),
            Paragraph(f"{photo['slide']}<br/><font size=7 color='#666666'>{photo['title']}</font>",
                      STYLES["cell"]),
            Paragraph(f"{photo['width_px']} × {photo['height_px']} px", STYLES["cell"]),
            Paragraph(build_prompt(photo, company, score), STYLES["cell"]),
        ])
    story.append(_table(rows, [17 * mm, 38 * mm, 23 * mm, 96 * mm]))

    return _build(story, "Image Replacement Brief", "Image Replacement Brief · Granuler")

import copy
import io
import re
import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

import yaml as _yaml

TEMPLATE_PATH = Path(__file__).parent.parent / "assets" / "granuler_template.pptx"
PILLAR_SLIDE_START = 43  # slide index 43 = slide number 44 (0-based)

_cfg_path = Path(__file__).parent / "config.yaml"
with open(_cfg_path) as _f:
    _cfg = _yaml.safe_load(_f)

PILLAR_COUNT: int = _cfg.get("pillar_count", 10)
SUBTOPICS_PER_PILLAR: int = _cfg.get("subtopics_per_pillar", 4)


_DEFAULT_PT = 14.0
_LINE_SPACING = 1.2
# Shrinking past this looks worse than the overflow it fixes: the text ends up
# visibly smaller than the identical card beside it, and PowerPoint does not
# clip overflow anyway, so a slight spill into the whitespace below reads
# better than 7pt type. Text needing more than this is left alone.
_MIN_FONT_SCALE = 0.8
_MIN_LABEL_SCALE = 0.6
# Beyond this a text frame is prose, not a card title.
_LABEL_WORDS = 6
_WIDTH_SAFETY = 1.06


def _strip_highlight(tf):
    """Drop any highlight fill carried over from the template.

    The template's overall-score run is highlighted yellow. Filling the run
    keeps its formatting, so every client's score shipped with a marker-pen
    block behind it.
    """
    for element in tf._txBody.iter():
        if element.tag == qn("a:highlight"):
            element.getparent().remove(element)


def _text_height_pt(tf, width_pt: float) -> float:
    """Estimate the height this text needs, wrapped to width_pt.

    Helvetica metrics stand in for the template's Arial - the two are close
    enough for a shrink factor, and reportlab is already a dependency.
    """
    from reportlab.pdfbase.pdfmetrics import stringWidth

    total = 0.0
    for para in tf.paragraphs:
        size, face = _para_font(para)
        words = "".join(r.text for r in para.runs).split()
        if not words:
            total += size * _LINE_SPACING
            continue
        lines, current = 1, ""
        for word in words:
            trial = f"{current} {word}".strip()
            if current and stringWidth(trial, face, size) * _WIDTH_SAFETY > width_pt:
                lines += 1
                current = word
            else:
                current = trial
        total += lines * size * _LINE_SPACING
    return total


def _para_font(para) -> tuple[float, str]:
    """The size and reportlab face to measure a paragraph with.

    Card titles are bold, and bold Arial is a few percent wider than regular.
    Measuring them with the regular face under-reads the width, which is enough
    to let a title wrap onto the description below it.
    """
    size = next((r.font.size.pt for r in para.runs if r.font.size), None) or _DEFAULT_PT
    bold = any(r.font.bold for r in para.runs)
    return size, "Helvetica-Bold" if bold else "Helvetica"


def _label_width_scale(tf, width_pt: float) -> float:
    """How much a short label must shrink to stay on one line."""
    from reportlab.pdfbase.pdfmetrics import stringWidth

    worst = 1.0
    for para in tf.paragraphs:
        size, face = _para_font(para)
        text = "".join(r.text for r in para.runs).strip()
        if not text:
            continue
        # PowerPoint's own metrics differ slightly from reportlab's, and a
        # near-miss still wraps, so leave a little headroom.
        needed = stringWidth(text, face, size) * _WIDTH_SAFETY
        if needed > width_pt:
            worst = min(worst, width_pt / needed)
    return worst


def _fit_text(shape):
    """Shrink text that would otherwise spill out of its box.

    python-pptx writes a bare <a:normAutofit/>, which PowerPoint ignores until
    someone edits the text - so generated decks rendered at full size and
    overflowed. The scale is therefore computed here and written explicitly,
    the way PowerPoint itself would.
    """
    tf = shape.text_frame
    width_pt = (shape.width - tf.margin_left - tf.margin_right) / 12700
    height_pt = (shape.height - tf.margin_top - tf.margin_bottom) / 12700
    if width_pt <= 0 or height_pt <= 0:
        return
    if len(tf.text.split()) <= _LABEL_WORDS:
        # A card title sits in a one-line box directly above its own
        # description, so wrapping does not spill into whitespace - it lands on
        # the text below. What matters here is keeping it on one line, not the
        # box height, which for these boxes is barely a line tall to begin with.
        scale, floor = _label_width_scale(tf, width_pt), _MIN_LABEL_SCALE
    else:
        needed = _text_height_pt(tf, width_pt)
        scale, floor = height_pt / needed, _MIN_FONT_SCALE
    if scale >= 1.0 or scale < floor:
        return
    body_pr = tf._txBody.find(qn("a:bodyPr"))
    if body_pr is None:
        return
    for existing in body_pr.findall(qn("a:normAutofit")):
        body_pr.remove(existing)
    autofit = body_pr.makeelement(qn("a:normAutofit"), {})
    autofit.set("fontScale", str(int(scale * 100000)))
    autofit.set("lnSpcReduction", "10000")
    body_pr.append(autofit)


def _set_text(shape, text: str):
    """Replace all text in a shape's text frame, preserving first run's formatting."""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    para = tf.paragraphs[0]
    ref_run = para.runs[0] if para.runs else None

    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)

    for r in para.runs[1:]:
        r._r.getparent().remove(r._r)

    if ref_run:
        ref_run.text = text
    else:
        para.text = text

    _strip_highlight(tf)
    _fit_text(shape)


def _set_bullet_list(shape, items: list[str]):
    """Set text frame to a bullet list, one paragraph per item."""
    from pptx.oxml.ns import qn
    from lxml import etree

    tf = shape.text_frame
    if not tf.paragraphs:
        return

    # Capture formatting from first paragraph/run
    first_para = tf.paragraphs[0]
    ref_run = first_para.runs[0] if first_para.runs else None

    # Remove all existing paragraphs except the first
    parent = first_para._p.getparent()
    existing = list(tf.paragraphs)
    for p in existing[1:]:
        parent.remove(p._p)

    # Set first item in existing first paragraph
    if items:
        if ref_run:
            ref_run.text = items[0]
            for r in first_para.runs[1:]:
                r._r.getparent().remove(r._r)
        else:
            first_para.text = items[0]

    # Add remaining items as new paragraphs cloned from first
    for item in items[1:]:
        new_p = copy.deepcopy(first_para._p)
        # Clear runs in cloned para and set text
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        # Create a run with text
        r_elem = copy.deepcopy(first_para._p.findall(qn("a:r"))[0]) if first_para._p.findall(qn("a:r")) else etree.SubElement(new_p, qn("a:r"))
        r_elem.find(qn("a:t")).text = item if r_elem.find(qn("a:t")) is not None else None
        if r_elem.find(qn("a:t")) is None:
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = item
        else:
            r_elem.find(qn("a:t")).text = item
        new_p.append(r_elem)
        parent.append(new_p)

    tf.word_wrap = True
    _strip_highlight(tf)
    _fit_text(shape)


def _replace_xyz(prs: Presentation, company_name: str):
    """Global replace 'XYZ' with company_name across all text runs."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "XYZ" in run.text:
                        run.text = run.text.replace("XYZ", company_name)


def _get_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _calc_maturity_band(score: float) -> str:
    if score < 40:
        return "At Risk Zone"
    elif score < 60:
        return "Developing Zone"
    elif score < 76:
        return "Managed Zone"
    elif score < 90:
        return "Advanced Zone"
    return "Leading Zone"


def _calc_pillar_score(subtopics: list[dict], subtopics_per_pillar: int = SUBTOPICS_PER_PILLAR) -> float:
    if not subtopics:
        return 0.0
    total = sum(s["score"] for s in subtopics)
    max_score = subtopics_per_pillar * 5
    return (total / max_score) * 10


def _calc_overall_score(pillars: list[dict], subtopics_per_pillar: int = SUBTOPICS_PER_PILLAR) -> float:
    pillar_scores = [_calc_pillar_score(p["subtopics"], subtopics_per_pillar) for p in pillars]
    return sum(pillar_scores) / len(pillar_scores) * 10 if pillar_scores else 0.0


def generate_report(
    intake: dict,
    pillars: list[dict],
    llm_global: dict,
    llm_pillars: list[dict],
    llm_narrative: dict | None = None,
    llm_context: dict | None = None,
    llm_architecture: dict | None = None,
    llm_findings: dict | None = None,
    llm_conditional: dict | None = None,
    llm_roadmap: dict | None = None,
    llm_closing: dict | None = None,
    llm_prior_work: dict | None = None,
) -> bytes:
    prs = Presentation(str(TEMPLATE_PATH))
    slides = prs.slides

    company = intake["company_name"]
    overall_score = _calc_overall_score(pillars)
    maturity_band = _calc_maturity_band(overall_score)

    # Global XYZ replacement
    _replace_xyz(prs, company)

    # --- Slide 1: Title ---
    s1 = slides[0]
    sh = _get_shape_by_name(s1, "Text 3")
    if sh:
        _set_text(sh, company)
    sh = _get_shape_by_name(s1, "Text 7")
    if sh:
        _set_text(sh, intake.get("assessment_date", ""))
    sh = _get_shape_by_name(s1, "Text 11")
    if sh:
        _set_text(sh, intake.get("assessor", "Ravi Kajaria"))

    # --- Slide 4: Maturity Summary ---
    s4 = slides[3]
    sh = _get_shape_by_name(s4, "Text 1")
    if sh:
        _set_text(sh, f"{overall_score:.1f}")
    sh = _get_shape_by_name(s4, "Text 7")
    if sh:
        _set_text(sh, f"Maturity Band: {maturity_band}")
    sh = _get_shape_by_name(s4, "Text 8")
    if sh:
        _set_text(sh, llm_global.get("maturity_summary", ""))
    sh = _get_shape_by_name(s4, "Text 10")
    if sh:
        _set_text(sh, llm_global.get("score_interpretation", ""))

    # --- Slide 8: Pillar Overview ---
    s8 = slides[7]
    sh = _get_shape_by_name(s8, "Text 3")
    if sh:
        _set_text(sh, f"Strongest Area — {llm_global.get('strongest_area', '')}")
    sh = _get_shape_by_name(s8, "Text 5")
    if sh:
        _set_text(sh, f"Weakest Areas — {llm_global.get('weakest_areas', '')}")

    # --- Slide 9: Risk Heatmap ---
    s9 = slides[8]
    sh = _get_shape_by_name(s9, "Text 4")
    if sh:
        _set_bullet_list(sh, llm_global.get("high_priority_risks", []))
    sh = _get_shape_by_name(s9, "Text 7")
    if sh:
        _set_bullet_list(sh, llm_global.get("high_impact_risks", []))
    sh = _get_shape_by_name(s9, "Text 10")
    if sh:
        _set_bullet_list(sh, llm_global.get("medium_risks", []))

    # --- Slide 25: 90-Day Plan ---
    s25 = slides[24]
    sh = _get_shape_by_name(s25, "Text 4")
    if sh:
        _set_bullet_list(sh, llm_global.get("days_1_30", []))
    sh = _get_shape_by_name(s25, "Text 6")
    if sh:
        _set_bullet_list(sh, llm_global.get("days_31_60", []))
    sh = _get_shape_by_name(s25, "Text 8")
    if sh:
        _set_bullet_list(sh, llm_global.get("days_61_90", []))

    # --- Slide 26: 12-Month Roadmap ---
    s26 = slides[25]
    sh = _get_shape_by_name(s26, "Text 6")
    if sh:
        _set_bullet_list(sh, llm_global.get("q1_items", []))
    sh = _get_shape_by_name(s26, "Text 11")
    if sh:
        _set_bullet_list(sh, llm_global.get("q2_items", []))
    sh = _get_shape_by_name(s26, "Text 16")
    if sh:
        _set_bullet_list(sh, llm_global.get("q3_items", []))
    sh = _get_shape_by_name(s26, "Text 21")
    if sh:
        _set_bullet_list(sh, llm_global.get("q4_items", []))

    # --- Slide 54: Thank You ---
    s54 = slides[53]
    sh = _get_shape_by_name(s54, "Text 2")
    if sh:
        _set_text(sh, llm_global.get("closing_message", f"We are committed to helping {company} transform technology from an operational tool into a strategic growth enabler."))

    # --- Slides 44-53: Pillar deep dives ---
    for idx, pillar_data in enumerate(pillars):
        slide_idx = PILLAR_SLIDE_START + idx
        if slide_idx >= len(slides):
            break
        sp = slides[slide_idx]
        pillar_score = _calc_pillar_score(pillar_data["subtopics"])
        llm = llm_pillars[idx] if idx < len(llm_pillars) else {}

        sh = _get_shape_by_name(sp, "Text 4")
        if sh:
            _set_text(sh, f"SCORE: {pillar_score:.1f} / 10")
        sh = _get_shape_by_name(sp, "Text 6")
        if sh:
            _set_text(sh, llm.get("observation", ""))
        sh = _get_shape_by_name(sp, "Text 8")
        if sh:
            _set_text(sh, llm.get("business_impact", ""))
        sh = _get_shape_by_name(sp, "Text 10")
        if sh:
            _set_text(sh, llm.get("rec1", ""))
        sh = _get_shape_by_name(sp, "Text 11")
        if sh:
            _set_text(sh, llm.get("rec2", ""))
        sh = _get_shape_by_name(sp, "Text 12")
        if sh:
            _set_text(sh, llm.get("rec3", ""))

    # --- Narrative slides (Phase 1 — dynamic for any client) ---
    if llm_narrative:
        nav = llm_narrative

        # Slide 7: Business Drivers
        s7 = slides[6]
        drivers = nav.get("business_drivers", [])
        for i, (title_name, desc_name) in enumerate([("Text 2","Text 3"),("Text 4","Text 5"),("Text 6","Text 7"),("Text 8","Text 9")]):
            if i < len(drivers):
                sh = _get_shape_by_name(s7, title_name)
                if sh: _set_text(sh, drivers[i].get("title", ""))
                sh = _get_shape_by_name(s7, desc_name)
                if sh: _set_text(sh, drivers[i].get("description", ""))

        # Slide 12: Weakest pillar spotlight
        s12 = slides[11]
        worst = min(pillars, key=lambda p: _calc_pillar_score(p["subtopics"]))
        worst_score = _calc_pillar_score(worst["subtopics"])
        sh = _get_shape_by_name(s12, "Text 0")
        if sh: _set_text(sh, f"{worst['pillar']} Gap")
        sh = _get_shape_by_name(s12, "Text 1")
        if sh: _set_text(sh, f"{worst_score:.1f}")
        sh = _get_shape_by_name(s12, "Text 3")
        if sh: _set_text(sh, f"{worst['pillar']} Score — the weakest pillar in the assessment")
        issues = nav.get("weakest_pillar_issues", [])
        for i, (t, d) in enumerate([("Text 6","Text 7"),("Text 9","Text 10"),("Text 12","Text 13")]):
            if i < len(issues):
                sh = _get_shape_by_name(s12, t)
                if sh: _set_text(sh, issues[i].get("title", ""))
                sh = _get_shape_by_name(s12, d)
                if sh: _set_text(sh, issues[i].get("description", ""))
        impacts = nav.get("weakest_pillar_impacts", [])
        for i, (t, d) in enumerate([("Text 16","Text 17"),("Text 19","Text 20")]):
            if i < len(impacts):
                sh = _get_shape_by_name(s12, t)
                if sh: _set_text(sh, impacts[i].get("emoji_title", ""))
                sh = _get_shape_by_name(s12, d)
                if sh: _set_text(sh, impacts[i].get("description", ""))

        # Slide 18: Quick Wins
        s18 = slides[17]
        qw = nav.get("quick_wins", [])
        for i, (t, d) in enumerate([("Text 6","Text 7"),("Text 12","Text 13"),("Text 18","Text 19"),("Text 24","Text 25"),("Text 30","Text 31"),("Text 36","Text 37")]):
            if i < len(qw):
                sh = _get_shape_by_name(s18, t)
                if sh: _set_text(sh, qw[i].get("title", ""))
                sh = _get_shape_by_name(s18, d)
                if sh: _set_text(sh, qw[i].get("description", ""))

        # Slide 24: Cost of Inaction
        s24 = slides[23]
        risks = nav.get("inaction_risks", [])
        for i, (t, d) in enumerate([("Text 3","Text 4"),("Text 6","Text 7"),("Text 9","Text 10"),("Text 12","Text 13")]):
            if i < len(risks):
                sh = _get_shape_by_name(s24, t)
                if sh: _set_text(sh, risks[i].get("emoji_title", ""))
                sh = _get_shape_by_name(s24, d)
                if sh: _set_text(sh, risks[i].get("description", ""))
        sh = _get_shape_by_name(s24, "Text 15")
        if sh: _set_text(sh, nav.get("inaction_closing", ""))

        # Slide 41: Expected Outcomes
        s41 = slides[40]
        outcomes = nav.get("expected_outcomes", [])
        for i, (t, d) in enumerate([("Text 2","Text 3"),("Text 4","Text 5"),("Text 6","Text 7"),("Text 8","Text 9")]):
            if i < len(outcomes):
                sh = _get_shape_by_name(s41, t)
                if sh: _set_text(sh, outcomes[i].get("title", ""))
                sh = _get_shape_by_name(s41, d)
                if sh: _set_text(sh, outcomes[i].get("description", ""))

    # --- Newly-dynamic slides ---
    # Slides with no content supplied are deleted rather than left carrying the
    # template's original client's text. An absent LLM block therefore means
    # "omit the slide", never "ship whatever the template said".
    to_delete: set[int] = set()

    if llm_context:
        _fill_context_slides(slides, llm_context, intake)
    if llm_architecture:
        _fill_architecture_slides(slides, llm_architecture)
    if llm_findings:
        _fill_findings_slides(slides, llm_findings)
    if llm_roadmap:
        _fill_roadmap_slides(slides, llm_roadmap)
    if llm_closing:
        _fill_closing_slides(slides, llm_closing)

    if llm_conditional:
        to_delete |= _fill_conditional_slides(slides, llm_conditional)
    else:
        to_delete |= set(CONDITIONAL_SLIDES)

    if llm_prior_work:
        _fill_prior_work_slides(slides, llm_prior_work)
    else:
        to_delete |= set(HISTORY_SLIDES)

    # Deletion must come after every fill: they are all keyed on absolute
    # slide indices. The guards run after it so they only inspect slides that
    # are actually shipping — a gated slide still holds template text until it
    # is removed, and flagging that would be a false positive.
    _delete_slides(prs, to_delete)
    _remove_foreign_images(prs, intake.get("company_name", ""), overall_score)
    _note_image_prompts(prs, intake.get("company_name", ""), overall_score)

    _assert_no_placeholder(prs)
    _assert_no_foreign_content(prs, intake, pillars)
    _assert_no_foreign_images(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Newly-dynamic slides.
#
# Every slide below carried hardcoded Uni-tech Automation content until
# 2026-09-01 and shipped verbatim to every client. Indices are 0-based
# (slide 1 = index 0). See api/slide_map.py for the full documented map.
#
# HISTORY_SLIDES describe work Granuler has already delivered for the client.
# CONDITIONAL_SLIDES are domain-dependent. Both are deleted rather than filled
# with invented content — see _delete_slides() and the ordering rule there.
# ---------------------------------------------------------------------------

S_HOOK, S_COMPANY, S_INTERPRET = 1, 2, 5
S_CORE_RISK, S_SECURITY, S_ARCH_COMPARE = 9, 10, 12
S_REPORTING, S_HR, S_INFRA, S_VENDOR = 13, 14, 15, 16
S_JOURNEY, S_ARCH_NOW, S_ARCH_FUTURE, S_QUALITY = 18, 19, 20, 22
S_DELIVERY, S_WHY_GRANULER, S_PATH_FORWARD, S_CORE_PROCESS = 26, 27, 28, 29
S_PROGRESS, S_RISK_MAP, S_WINS, S_VALUE, S_INACTION2 = 30, 31, 32, 33, 34
S_TOP10, S_ROADMAP_PHASES, S_TIMELINE = 35, 36, 37
S_ACT_NOW, S_CLOSING_STATS = 39, 41

HISTORY_SLIDES = (S_PROGRESS, S_WINS, S_VALUE)
CONDITIONAL_SLIDES = (S_CORE_RISK, S_HR, S_VENDOR, S_QUALITY, S_CORE_PROCESS)

_PAIRS_4A = [("Text 4", "Text 5"), ("Text 8", "Text 9"), ("Text 12", "Text 13"), ("Text 16", "Text 17")]
_PAIRS_4B = [("Text 3", "Text 4"), ("Text 6", "Text 7"), ("Text 9", "Text 10"), ("Text 12", "Text 13")]
_PAIRS_3A = [("Text 4", "Text 5"), ("Text 8", "Text 9"), ("Text 12", "Text 13")]
_PAIRS_3B = [("Text 1", "Text 2"), ("Text 3", "Text 4"), ("Text 5", "Text 6")]
_PAIRS_3C = [("Text 10", "Text 11"), ("Text 13", "Text 14"), ("Text 16", "Text 17")]
_PAIRS_3D = [("Text 4", "Text 5"), ("Text 6", "Text 7"), ("Text 8", "Text 9")]


def _fill(slide, shape_name: str, text: str):
    """Write text to a named shape, blanking it when there is nothing to write.

    Blank must mean "clear", not "keep". Only the newly-dynamic slides use this,
    and on those the template text belongs to a different client — leaving it in
    place is the bug this module was reworked to remove. Slide 3's "Industries
    Served" shipped Uni-tech's industries this way during the 2026-09-01 smoke
    test, because the field came back empty.
    """
    shape = _get_shape_by_name(slide, shape_name)
    if shape is not None:
        _set_text(shape, text or "")


def _clear(slide, *shape_names):
    """Blank a slot the LLM returned no content for.

    Leaving it alone would ship the template's original client's text, which is
    the entire bug this module was reworked to prevent. An empty slot is the
    lesser failure and is what the caller gets when a list comes back short.
    """
    for shape_name in shape_names:
        shape = _get_shape_by_name(slide, shape_name)
        if shape is not None:
            _set_text(shape, "")


def _fill_pairs(slide, pairs, items, tkey="title", dkey="description"):
    items = items or []
    for index, (title_shape, desc_shape) in enumerate(pairs):
        item = items[index] if index < len(items) else None
        if not isinstance(item, dict):
            _clear(slide, title_shape, desc_shape)
            continue
        _fill(slide, title_shape, item.get(tkey, ""))
        _fill(slide, desc_shape, item.get(dkey, ""))


def _fill_stats(slide, triples, items):
    items = items or []
    for index, (val_shape, label_shape, desc_shape) in enumerate(triples):
        item = items[index] if index < len(items) else None
        if not isinstance(item, dict):
            _clear(slide, val_shape, label_shape, desc_shape)
            continue
        _fill(slide, val_shape, str(item.get("value", "")))
        _fill(slide, label_shape, item.get("label", ""))
        _fill(slide, desc_shape, item.get("description", ""))


def _fill_context_slides(slides, ctx: dict, intake: dict):
    """Slides 2, 3, 6, 27, 29."""
    s = slides[S_HOOK]
    _fill(s, "Text 0", ctx.get("hook_question", ""))
    _fill(s, "Text 1", ctx.get("growth_framing", ""))
    pillar_labels = ctx.get("growth_pillars") or []
    for index, shape_name in enumerate(["Text 4", "Text 7", "Text 10", "Text 13"]):
        _fill(s, shape_name, pillar_labels[index] if index < len(pillar_labels) else "")
    _fill(s, "Text 15", ctx.get("strategic_shift", ""))

    s = slides[S_COMPANY]
    _fill(s, "Text 1", ctx.get("company_description", ""))
    _fill(s, "Text 2", ctx.get("expansion_note", ""))
    _fill(s, "Text 5", ctx.get("products_line", ""))
    _fill(s, "Text 8", ctx.get("industries_line", ""))

    _fill(slides[S_INTERPRET], "Text 1", ctx.get("score_interpretation_long", ""))

    # Title is derived, not generated — it is a plain fact about the engagement.
    s = slides[S_DELIVERY]
    base = (intake.get("granuler_location") or "Mumbai").strip()
    client_loc = (intake.get("locations") or "").split(",")[0].strip()
    if client_loc and client_loc.lower() != base.lower():
        _fill(s, "Text 0", f"Delivery Model: {base} → {client_loc}")
    else:
        _fill(s, "Text 0", f"Delivery Model: {base}" if base else "Delivery Model")
    _fill(s, "Text 1", ctx.get("delivery_description", ""))
    _fill(s, "Text 3", ctx.get("delivery_note", ""))
    _fill_pairs(s, _PAIRS_3D, ctx.get("delivery_modes", []))

    s = slides[S_PATH_FORWARD]
    _fill(s, "Text 1", ctx.get("path_forward_intro", ""))
    _fill_pairs(s, _PAIRS_3A, ctx.get("path_forward_items", []))
    _fill(s, "Text 14", ctx.get("path_forward_closing", ""))


def _fill_architecture_slides(slides, arch: dict):
    """Slides 13, 19, 20, 21."""
    s = slides[S_ARCH_COMPARE]
    _fill_pairs(s, _PAIRS_4A, arch.get("current_arch", []))
    _fill_pairs(
        s,
        [("Text 20", "Text 21"), ("Text 23", "Text 24"), ("Text 26", "Text 27"), ("Text 29", "Text 30")],
        arch.get("future_arch", []),
    )

    s = slides[S_JOURNEY]
    _fill(s, "Text 1", arch.get("journey_intro", ""))
    stage_pairs = [("Text 2", "Text 3"), ("Text 4", "Text 5"), ("Text 6", "Text 7"), ("Text 8", "Text 9")]
    # The template's stage titles carry traffic-light emoji; keep them.
    dots = ["\U0001F534", "\U0001F7E0", "\U0001F7E1", "\U0001F7E2"]
    stages = [
        {"title": f"{dot} {st.get('title', '')}".strip(), "description": st.get("description", "")}
        for dot, st in zip(dots, arch.get("journey_stages", []) or [])
        if isinstance(st, dict)
    ]
    _fill_pairs(s, stage_pairs, stages)

    s = slides[S_ARCH_NOW]
    _fill_pairs(s, _PAIRS_3B, arch.get("current_layers", []))
    _fill(s, "Text 8", arch.get("current_summary", ""))
    _fill_pairs(s, _PAIRS_3C, arch.get("current_risks", []))

    s = slides[S_ARCH_FUTURE]
    _fill_pairs(s, _PAIRS_3B, arch.get("future_layers", []))
    _fill(s, "Text 8", arch.get("future_summary", ""))
    _fill_pairs(s, _PAIRS_3C, arch.get("future_gains", []))


def _fill_findings_slides(slides, find: dict):
    """Slides 11, 14, 16."""
    s = slides[S_SECURITY]
    _fill(s, "Text 1", find.get("security_intro", ""))
    _fill(s, "Text 3", find.get("security_note", ""))
    _fill_pairs(
        s,
        [("Text 6", "Text 7"), ("Text 10", "Text 11"), ("Text 14", "Text 15"), ("Text 18", "Text 19")],
        find.get("security_findings", []),
    )

    s = slides[S_REPORTING]
    _fill_pairs(s, _PAIRS_3B, find.get("reporting_flow", []))
    _fill_pairs(
        s,
        [("Text 10", "Text 11"), ("Text 14", "Text 15"), ("Text 18", "Text 19")],
        find.get("reporting_current", []),
    )
    _fill(s, "Text 21", find.get("reporting_recommendation", ""))

    s = slides[S_INFRA]
    _fill(s, "Text 1", find.get("infra_intro", ""))
    _fill_pairs(s, _PAIRS_4A, find.get("infra_findings", []))
    _fill(s, "Text 18", find.get("infra_closing", ""))


def _fill_conditional_slides(slides, cond: dict) -> set[int]:
    """Slides 10, 15, 17, 23, 30. Returns the indices to delete."""
    drop = set()

    block = cond.get("core_system_risk") or {}
    if block.get("applicable"):
        s = slides[S_CORE_RISK]
        _fill(s, "Text 0", block.get("title", ""))
        _fill(s, "Text 2", block.get("warning", ""))
        _fill_pairs(s, [("Text 5", "Text 6"), ("Text 9", "Text 10"), ("Text 13", "Text 14")], block.get("impacts", []))
        _fill(s, "Text 15", block.get("closing", ""))
    else:
        drop.add(S_CORE_RISK)

    block = cond.get("hr_opportunity") or {}
    if block.get("applicable"):
        s = slides[S_HR]
        _fill(s, "Text 1", block.get("intro", ""))
        _fill_pairs(s, [("Text 2", "Text 3"), ("Text 4", "Text 5"), ("Text 6", "Text 7")], block.get("items", []))
    else:
        drop.add(S_HR)

    block = cond.get("vendor_governance") or {}
    if block.get("applicable"):
        s = slides[S_VENDOR]
        _fill(s, "Text 0", block.get("title", ""))
        _fill_pairs(s, _PAIRS_3A, block.get("observations", []))
        _fill(s, "Text 15", block.get("action_taken", ""))
    else:
        drop.add(S_VENDOR)

    block = cond.get("quality_process") or {}
    if block.get("applicable"):
        s = slides[S_QUALITY]
        _fill(s, "Text 1", block.get("intro", ""))
        _fill_pairs(s, [("Text 4", "Text 5"), ("Text 7", "Text 8"), ("Text 10", "Text 11")], block.get("within_systems", []))
        _fill_pairs(s, [("Text 14", "Text 15"), ("Text 17", "Text 18"), ("Text 20", "Text 21")], block.get("outside_systems", []))
    else:
        drop.add(S_QUALITY)

    block = cond.get("core_process_observations") or {}
    if block.get("applicable"):
        s = slides[S_CORE_PROCESS]
        _fill(s, "Text 0", block.get("title", ""))
        _fill(s, "Text 1", block.get("intro", ""))
        _fill_pairs(s, _PAIRS_4B, block.get("findings", []))
    else:
        drop.add(S_CORE_PROCESS)

    return drop


def _fill_roadmap_slides(slides, road: dict):
    """Slides 32, 36, 37, 38."""
    s = slides[S_RISK_MAP]
    _fill(s, "Text 1", road.get("risk_mapping_intro", ""))
    _fill_pairs(
        s,
        [("Text 7", "Text 8"), ("Text 10", "Text 11"), ("Text 13", "Text 14"), ("Text 16", "Text 17"), ("Text 19", "Text 20")],
        road.get("risk_mapping", []),
        tkey="risk",
        dkey="initiative",
    )

    _fill_pairs(
        slides[S_TOP10],
        [("Text 3", "Text 4"), ("Text 7", "Text 8"), ("Text 11", "Text 12"), ("Text 15", "Text 16"),
         ("Text 19", "Text 20"), ("Text 23", "Text 24"), ("Text 27", "Text 28"), ("Text 31", "Text 32"),
         ("Text 35", "Text 36"), ("Text 39", "Text 40")],
        road.get("top_priorities", []),
    )

    s = slides[S_ROADMAP_PHASES]
    _fill_pairs(s, _PAIRS_3B, road.get("roadmap_phases", []))
    _fill(s, "Text 7", road.get("roadmap_closing", ""))

    _fill_pairs(
        slides[S_TIMELINE],
        [("Text 5", "Text 6"), ("Text 10", "Text 11"), ("Text 15", "Text 16"), ("Text 20", "Text 21")],
        road.get("timeline_quarters", []),
    )


def _fill_closing_slides(slides, close: dict):
    """Slides 28, 35, 40, 42."""
    s = slides[S_WHY_GRANULER]
    _fill(s, "Text 1", close.get("why_granuler_intro", ""))
    _fill_pairs(
        s,
        [("Text 4", "Text 5"), ("Text 8", "Text 9"), ("Text 12", "Text 13"), ("Text 16", "Text 17"), ("Text 20", "Text 21")],
        close.get("why_granuler_items", []),
    )

    s = slides[S_INACTION2]
    _fill(s, "Text 1", close.get("inaction_intro", ""))
    _fill_pairs(s, _PAIRS_4B, close.get("inaction_items", []))
    _fill(s, "Text 15", close.get("inaction_principle", ""))

    s = slides[S_ACT_NOW]
    _fill(s, "Text 1", close.get("act_now_intro", ""))
    _fill_pairs(s, _PAIRS_4B, close.get("act_now_items", []))

    s = slides[S_CLOSING_STATS]
    _fill_stats(
        s,
        [("Text 1", "Text 2", "Text 3"), ("Text 4", "Text 5", "Text 6"), ("Text 7", "Text 8", "Text 9")],
        close.get("closing_stats", []),
    )
    _fill(s, "Text 10", close.get("closing_statement", ""))


def _fill_prior_work_slides(slides, prior: dict):
    """Slides 31, 33, 34. Only called when the client has prior delivered work."""
    s = slides[S_PROGRESS]
    _fill(s, "Text 1", prior.get("progress_intro", ""))
    _fill_pairs(
        s,
        [("Text 4", "Text 5"), ("Text 8", "Text 9"), ("Text 12", "Text 13"), ("Text 17", "Text 18")],
        prior.get("progress_items", []),
    )

    s = slides[S_WINS]
    _fill_pairs(s, [("Text 3", "Text 4"), ("Text 6", "Text 7")], prior.get("governance_wins", []))
    _fill_pairs(s, [("Text 10", "Text 11"), ("Text 13", "Text 14")], prior.get("operational_wins", []))

    s = slides[S_VALUE]
    _fill(s, "Text 1", prior.get("value_intro", ""))
    _fill_stats(
        s,
        [("Text 2", "Text 3", "Text 4"), ("Text 5", "Text 6", "Text 7"), ("Text 8", "Text 9", "Text 10")],
        prior.get("value_stats", []),
    )


def _delete_slides(prs: Presentation, indices) -> None:
    """Remove slides by 0-based index.

    Descending order is mandatory: removing a low index first would shift every
    higher one. Nothing may write to a slide by index after this runs, so
    generate_report() calls it exactly once, immediately before save.
    """
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    for idx in sorted(indices, reverse=True):
        if 0 <= idx < len(entries):
            id_list.remove(entries[idx])


def _assert_no_placeholder(prs: Presentation, placeholder: str = "the client company") -> None:
    """Fail loudly rather than ship an unresolved LLM placeholder.

    MISTAKES.md 2026-08-24: a case-sensitive restore let "The client company"
    reach 17 slides of a delivered deck. This is the backstop for that.
    """
    needle = placeholder.lower()
    hits = []
    for number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame and needle in shape.text_frame.text.lower():
                hits.append(f"slide {number} / {shape.name}")
    if hits:
        raise ValueError(
            f"Unresolved {placeholder!r} placeholder left in the deck: {', '.join(hits)}"
        )


# Fingerprints of the original hand-built Uni-tech Automation deck that the
# template was cut from. Any of these surviving into a finished deck means a
# slide was not filled — unless the client's own input mentions the same thing,
# which is why the check is made against the input corpus rather than absolutely.
# Pictures that carry another client's identity in their pixels. The text
# guards cannot see inside a raster, so these are matched by content hash and
# removed outright. The maturity-band staircase on slide 6 is captioned
# "(Uni-Tech)" and highlights that client's band, so it shipped a different
# company's name and score to every reader.
_FOREIGN_IMAGE_SHA1 = {
    "f6a948972937531bbc043358b66dccc5977578f0": "slide 6 maturity staircase captioned (Uni-Tech)",
}


def _remove_foreign_images(prs: Presentation, company: str = "",
                           overall_score: float | None = None) -> list[str]:
    """Delete pictures whose pixels name another client.

    A labelled placeholder is left in the gap, because otherwise the slide just
    looks blank and nothing says a replacement is owed.
    """
    removed = []
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.__class__.__name__ != "Picture":
                continue
            try:
                sha1 = shape.image.sha1
            except ValueError:
                continue
            if sha1 not in _FOREIGN_IMAGE_SHA1:
                continue
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            _add_image_placeholder(slide, left, top, width, height)
            _note_replacement(slide, width, height, company, overall_score)
            removed.append(_FOREIGN_IMAGE_SHA1[sha1])
    return removed


def _add_image_placeholder(slide, left, top, width, height) -> None:
    """A visible, client-safe marker where a replacement image belongs."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = "Image Placeholder"
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF2, 0xF8, 0xFA)
    line = box.line
    line.color.rgb = RGBColor(0x9F, 0xC4, 0xCC)
    line.width = Pt(1)

    frame = box.text_frame
    frame.word_wrap = True
    px_w = round(width / 914400 * 150)
    px_h = round(height / 914400 * 150)
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = (
        f"Image goes here — {px_w} x {px_h} px\n"
        "The prompt for it is in this slide's speaker notes."
    )
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x14, 0x60, 0x6F)


def _note_replacement(slide, width, height, company: str, overall_score: float | None) -> None:
    """The staircase prompt, in the notes of the slide it belongs to.

    The placeholder is a text box rather than a picture, so the photo collector
    does not see it and it would otherwise be the one slide with a gap and no
    prompt attached.
    """
    from .image_brief import build_prompt

    photo = {
        "slide": 6,  # the template slide this graphic belongs to
        "title": "How to Interpret the Score",
        "width_in": width / 914400,
        "height_in": height / 914400,
    }
    notes = slide.notes_slide.notes_text_frame
    existing = notes.text.strip()
    lines = [
        "REPLACE THIS PICTURE - it named a different client and was removed.",
        f"Render at {round(width / 914400 * 150)} x {round(height / 914400 * 150)} px.",
        "",
        "Prompt:",
        build_prompt(photo, company, overall_score),
    ]
    if existing:
        lines += ["", "---", existing]
    notes.text = "\n".join(lines)


def _note_image_prompts(prs: Presentation, company: str, overall_score: float) -> int:
    """Write each replaceable photograph's prompt into its own speaker notes.

    Pointing at these slides from an outside document does not work: the deck
    drops gated slides so its numbering is not the template's, and several
    titles are rewritten per client. Notes travel inside the file, so they
    cannot drift, and a viewer never sees them.
    """
    from .image_brief import build_prompt, photos_in

    photos = photos_in(prs)
    slides = list(prs.slides)
    for photo in photos:
        notes = slides[photo["slide"] - 1].notes_slide.notes_text_frame
        existing = notes.text.strip()
        notes.text = (
            f"IMAGE {photo['number']} OF {len(photos)} — REPLACE THIS PICTURE\n"
            f"Render at {photo['width_px']} x {photo['height_px']} px.\n\n"
            f"Prompt:\n{build_prompt(photo, company, overall_score)}\n"
            + (f"\n---\n{existing}" if existing else "")
        )
    return len(photos)


def _assert_no_foreign_images(prs: Presentation) -> None:
    """No picture naming another client may survive to the saved deck."""
    for number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.__class__.__name__ != "Picture":
                continue
            try:
                sha1 = shape.image.sha1
            except ValueError:
                continue
            if sha1 in _FOREIGN_IMAGE_SHA1:
                raise AssertionError(
                    f"slide {number} still carries {_FOREIGN_IMAGE_SHA1[sha1]}"
                )


_TEMPLATE_FINGERPRINTS = (
    "sap", "s/4hana", "1709", "fiori", "sydler", "uni-tech", "unitech",
    "windows 7", "windows 11", "pune", "mrp", "bom", "nas box",
    "₹10l", "₹10,00,000", "pcc", "mcc", "apfc",
    # Uni-tech's products and markets, from slides 2 and 3.
    "wiring harness", "automation manufacturing", "power generation",
    "engineering infrastructure", "industrial manufacturing",
)


def _input_corpus(intake: dict, pillars: list[dict]) -> str:
    parts = [str(v) for v in intake.values() if v]
    for pillar in pillars:
        parts.append(pillar.get("pillar", ""))
        for sub in pillar.get("subtopics", []):
            parts.extend([sub.get("subtopic", ""), sub.get("current_state_notes", ""), sub.get("evidence", "")])
    return " ".join(parts).lower()


def _mentions(haystack: str, term: str) -> bool:
    """Whole-token match. A plain substring test reports 'sap' inside 'whatsapp'."""
    return re.search(rf"(?<![a-z0-9]){re.escape(term)}(?![a-z0-9])", haystack) is not None


def _assert_no_foreign_content(prs: Presentation, intake: dict, pillars: list[dict]) -> None:
    """Fail loudly if another client's facts survived into this deck.

    The reported 2026-08-31 bug: 33 of 54 slides were never written to, so a
    deck for Nihaar Equipments shipped with SAP, Sydler, Windows 7 and Pune.
    A term is only a leak if it is NOT in this client's own input.
    """
    corpus = _input_corpus(intake, pillars)
    leaked: dict[str, list[int]] = {}
    for number, slide in enumerate(prs.slides, start=1):
        text = " ".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        ).lower()
        for term in _TEMPLATE_FINGERPRINTS:
            if _mentions(text, term) and not _mentions(corpus, term):
                leaked.setdefault(term, []).append(number)
    if leaked:
        detail = "; ".join(f"{term!r} on slides {nums}" for term, nums in sorted(leaked.items()))
        raise ValueError(f"Template content from another client survived into the deck: {detail}")

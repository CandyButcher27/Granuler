"""The stock-photo replacement brief."""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from api.image_brief import build_prompt, collect_photos, image_brief_pdf  # noqa: E402


def test_finds_the_stock_photographs_and_ignores_the_icons():
    """The template carries ~117 pictures; only the photographs need replacing.

    The rest are icons, rules and SVG glyphs that carry no client identity.
    """
    photos = collect_photos()
    assert 10 <= len(photos) <= 30, f"{len(photos)} photos - the filter has drifted"
    assert [p["number"] for p in photos] == list(range(1, len(photos) + 1))
    for photo in photos:
        assert photo["slide"] >= 1
        assert photo["width_px"] > 0 and photo["height_px"] > 0


def test_every_photo_gets_a_usable_prompt():
    from api.image_brief import SPECIAL_PROMPTS

    for photo in collect_photos():
        prompt = build_prompt(photo)
        if photo["slide"] in SPECIAL_PROMPTS:
            continue  # a diagram, briefed separately below
        assert "No text, no logos" in prompt
        assert prompt.rstrip().endswith("crop.")


def test_orientation_follows_the_shape():
    wide = build_prompt({"title": "x", "width_in": 10.0, "height_in": 5.0})
    tall = build_prompt({"title": "x", "width_in": 5.0, "height_in": 10.0})
    square = build_prompt({"title": "x", "width_in": 6.0, "height_in": 6.0})
    assert "wide landscape" in wide
    assert "tall portrait" in tall
    assert "square" in square


def test_renders_a_pdf():
    assert image_brief_pdf().startswith(b"%PDF-")


def test_renders_when_the_template_has_no_photographs():
    assert image_brief_pdf([]).startswith(b"%PDF-")


def test_the_maturity_diagram_names_this_client_not_the_template_s():
    """Slide 6's picture is a staircase captioned "(Uni-Tech)" in its pixels.

    A photo prompt is wrong for it twice over: it is a diagram, and its
    replacement has to carry the current client's name and band.
    """
    photo = next(p for p in collect_photos() if p["slide"] == 6)
    prompt = build_prompt(photo, "Nihaar Equipments", 41.5)
    assert "Nihaar Equipments" in prompt
    assert "41-60" in prompt
    assert "Uni-Tech" not in prompt
    assert "photograph" not in prompt.lower()


def test_the_band_highlighted_follows_the_score():
    from api.image_brief import band_range_for

    assert band_range_for(12.0) == "0-25"
    assert band_range_for(35.5) == "26-40"
    assert band_range_for(41.5) == "41-60"
    assert band_range_for(61.5) == "61-75"
    assert band_range_for(80.0) == "76-90"
    assert band_range_for(95.0) == "91-100"


def test_the_leaking_picture_never_reaches_a_generated_deck():
    from api.pptx_generator import _FOREIGN_IMAGE_SHA1, _remove_foreign_images
    from pptx import Presentation

    from api.image_brief import TEMPLATE_PATH

    prs = Presentation(str(TEMPLATE_PATH))
    assert _remove_foreign_images(prs), "the known foreign picture was not found to remove"
    surviving = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.__class__.__name__ != "Picture":
                continue
            try:
                surviving.add(shape.image.sha1)
            except ValueError:
                pass
    assert not (surviving & set(_FOREIGN_IMAGE_SHA1))

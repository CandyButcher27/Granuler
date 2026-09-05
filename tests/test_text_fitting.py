"""Text that would spill out of its box is shrunk before it collides.

python-pptx writes a bare <a:normAutofit/>, which PowerPoint ignores until the
text is edited by hand, so generated decks rendered at full size. The scale is
computed here instead, and the card titles are the case that matters: they sit
in one-line boxes directly above their own description, so a title that wraps
lands on the text below rather than in whitespace.
"""
import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

sys.path.insert(0, str(Path(__file__).parent.parent))

from api.pptx_generator import _fit_text, _set_text, _strip_highlight  # noqa: E402


def _textbox(width_in, height_in, text, size_pt=14):
    slide = Presentation().slides.add_slide(Presentation().slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(int(width_in * 914400)),
                                   Emu(int(height_in * 914400)))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    return box


def _font_scale(shape):
    autofit = shape.text_frame._txBody.find(qn("a:bodyPr")).find(qn("a:normAutofit"))
    return None if autofit is None else int(autofit.get("fontScale")) / 100000


def test_text_that_already_fits_is_left_alone():
    box = _textbox(4.0, 1.0, "Short label")
    _fit_text(box)
    assert _font_scale(box) is None


def test_a_wrapping_card_title_is_shrunk_rather_than_left_to_collide():
    """The exact slide-18 failure: a 3-word title wrapping onto its description."""
    box = _textbox(1.9, 0.24, "Implement Basic Automation")
    _fit_text(box)
    scale = _font_scale(box)
    assert scale is not None, "title was left to overlap the text beneath it"
    assert 0.6 <= scale < 1.0


def test_a_paragraph_is_never_shrunk_past_four_fifths():
    """Below that it reads as smaller than the identical card beside it."""
    box = _textbox(2.0, 0.3, " ".join(["word"] * 60))
    _fit_text(box)
    scale = _font_scale(box)
    assert scale is None or scale >= 0.8


def test_the_template_highlight_never_survives_a_fill():
    """The template's overall-score run is highlighted yellow."""
    box = _textbox(3.0, 1.0, "61.5")
    run = box.text_frame.paragraphs[0].runs[0]
    highlight = run._r.get_or_add_rPr().makeelement(qn("a:highlight"), {})
    run._r.get_or_add_rPr().append(highlight)
    assert b"highlight" in box.text_frame._txBody.xml.encode()

    _set_text(box, "35.5")
    assert "highlight" not in box.text_frame._txBody.xml

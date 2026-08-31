"""Locks the 2026-08-31 bug: a deck for one client shipped another client's facts.

Root cause was that 33 of 54 template slides were never written to, so the
original Uni-tech Automation prose reached every client. These tests assert the
deck a client actually receives contains nothing from the template's origin.
"""
import re
import sys
from pathlib import Path

import pytest
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent.parent))

from api.pptx_generator import (  # noqa: E402
    CONDITIONAL_SLIDES,
    HISTORY_SLIDES,
    generate_report,
)
from tests import fixtures as fx  # noqa: E402

FULL_CONTENT = dict(
    llm_global=fx.LLM_GLOBAL,
    llm_pillars=fx.LLM_PILLARS,
    llm_narrative=fx.LLM_NARRATIVE,
    llm_context=fx.LLM_CONTEXT,
    llm_architecture=fx.LLM_ARCHITECTURE,
    llm_findings=fx.LLM_FINDINGS,
    llm_conditional=fx.LLM_CONDITIONAL,
    llm_roadmap=fx.LLM_ROADMAP,
    llm_closing=fx.LLM_CLOSING,
)

# Every one of these appeared in the deck the client reported.
LEAKED_TERMS = [
    "SAP", "S/4HANA", "1709", "FIORI", "Sydler", "Uni-tech",
    "Windows 7", "Windows 11", "Pune", "MRP", "BOM",
    "₹10L", "₹10,00,000", "PCC", "MCC", "APFC",
]


def _deck(**overrides):
    kwargs = {**FULL_CONTENT, **overrides}
    data = generate_report(intake=fx.INTAKE, pillars=fx.PILLARS_RAW, **kwargs)
    return Presentation(__import__("io").BytesIO(data))


def _all_text(prs):
    return "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )


@pytest.fixture(scope="module")
def deck():
    return _deck()


@pytest.mark.parametrize("term", LEAKED_TERMS)
def test_no_leaked_term(deck, term):
    # Whole-token match: a substring test finds "sap" inside "whatsapp", which
    # is legitimately one of this client's systems.
    pattern = rf"(?<![a-z0-9]){re.escape(term.lower())}(?![a-z0-9])"
    assert not re.search(pattern, _all_text(deck).lower()), f"{term!r} leaked into the deck"


def test_no_unresolved_llm_placeholder(deck):
    assert "the client company" not in _all_text(deck).lower()


def test_client_facts_are_present(deck):
    text = _all_text(deck)
    for expected in ["Nihaar Equipments", "Mumbai", "Tally", "Stability Chambers"]:
        assert expected in text, f"{expected!r} missing from the deck"


def test_history_slides_dropped_without_prior_work(deck):
    # 3 history slides + the 3 conditional slides the fixture marks inapplicable.
    inapplicable = sum(
        1 for block in fx.LLM_CONDITIONAL.values() if not block.get("applicable")
    )
    assert len(deck.slides._sldIdLst) == 54 - len(HISTORY_SLIDES) - inapplicable


def test_prior_work_keeps_history_slides():
    prior = {
        "progress_intro": "Two improvements have already been delivered.",
        "progress_items": [{"title": "Contract Register", "description": "Contracts moved out of Excel."}],
        "governance_wins": [{"title": "Vendor Framework", "description": "A structured evaluation process is in place."}],
        "operational_wins": [{"title": "Backup Verified", "description": "Restore tested successfully."}],
        "value_intro": "Value was delivered before roadmap execution began.",
        "value_stats": [{"value": "10", "label": "Pillars Assessed", "description": "Baseline established"}],
    }
    with_history = _deck(llm_prior_work=prior)
    without = _deck()
    assert len(with_history.slides._sldIdLst) == len(without.slides._sldIdLst) + len(HISTORY_SLIDES)


def test_all_conditional_slides_dropped_when_none_apply():
    none_apply = {
        key: {**block, "applicable": False} for key, block in fx.LLM_CONDITIONAL.items()
    }
    prs = _deck(llm_conditional=none_apply)
    assert len(prs.slides._sldIdLst) == 54 - len(HISTORY_SLIDES) - len(CONDITIONAL_SLIDES)


def test_unfilled_slides_raise_rather_than_leak():
    """The guard is what makes a missing LLM block safe instead of silent."""
    with pytest.raises(ValueError, match="survived into the deck"):
        generate_report(
            intake=fx.INTAKE,
            pillars=fx.PILLARS_RAW,
            llm_global=fx.LLM_GLOBAL,
            llm_pillars=fx.LLM_PILLARS,
            llm_narrative=fx.LLM_NARRATIVE,
        )


def test_placeholder_left_in_content_raises():
    broken = {**fx.LLM_CONTEXT, "company_description": "The client company operates from Mumbai."}
    with pytest.raises(ValueError, match="placeholder"):
        _deck(llm_context=broken)

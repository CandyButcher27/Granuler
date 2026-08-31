"""Live end-to-end smoke test. Makes real LLM calls and costs real money.

Run:  .venv\\Scripts\\python tests\\smoke_live.py "Nihaar equipments Notes.txt"

Not collected by pytest (no test_ prefix) - it is the manual check to run before
declaring report changes done. The unit tests use fixtures, so they cannot catch
a prompt that returns an empty field or an inverted score; this can, and did:
three defects on 2026-09-01 (missing locations, inverted polarity on
negatively-named subtopics, and slide 3 keeping the template's industries).
"""
import io
import json
import os
import re
import sys
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

ROOT = Path(__file__).parent.parent

for line in (ROOT / ".env").read_text(encoding="utf-8").splitlines():
    if line.strip() and not line.startswith("#") and "=" in line:
        key, value = line.split("=", 1)
        os.environ.setdefault(key.strip(), value.strip().strip("\"'"))

from fastapi.testclient import TestClient  # noqa: E402
from pptx import Presentation  # noqa: E402

from api.main import app  # noqa: E402
from api.pptx_generator import _TEMPLATE_FINGERPRINTS  # noqa: E402

FORBIDDEN = _TEMPLATE_FINGERPRINTS + ("the client company",)


def main(notes_path: str, company: str) -> int:
    client = TestClient(app)
    notes = Path(notes_path).read_text(encoding="utf-8")

    start = time.time()
    resp = client.post("/extract-from-notes", json={"notes": notes, "company_name": company})
    print(f"extract-from-notes  {resp.status_code}  {time.time() - start:.1f}s")
    if resp.status_code != 200:
        print(resp.text[:2000])
        return 1

    extracted = resp.json()
    filled = sum(1 for v in extracted["intake"].values() if v)
    print(f"  intake filled: {filled}/{len(extracted['intake'])}")
    print(f"  pillars: {len(extracted['pillars'])} "
          f"subtopics: {[len(p['subtopics']) for p in extracted['pillars']]}")
    for field in ("locations", "core_systems", "industries_served"):
        print(f"  {field}: {extracted['intake'].get(field) or '(EMPTY)'}")

    payload = {
        "company_name": company,
        "assessor": "Ravi Kajaria",
        **{k: v for k, v in extracted["intake"].items() if v},
        "pillars": [
            {
                "pillar": p["pillar"],
                "subtopics": [
                    {
                        "subtopic": s["subtopic"],
                        "score": int(s["score"]),
                        "impact": s.get("impact", "Medium"),
                        "priority": s.get("priority", "Medium"),
                        "current_state_notes": s.get("current_state_notes", ""),
                    }
                    for s in p["subtopics"]
                ],
            }
            for p in extracted["pillars"]
        ],
    }

    start = time.time()
    resp = client.post("/generate-report", json=payload)
    print(f"generate-report     {resp.status_code}  {time.time() - start:.1f}s  {len(resp.content)} bytes")
    if resp.status_code != 200:
        print(resp.text[:3000])
        return 1

    out_dir = ROOT / "outputs"
    out_dir.mkdir(exist_ok=True)
    stem = company.replace(" ", "_")
    (out_dir / f"SMOKE_{stem}.pptx").write_bytes(resp.content)
    (out_dir / f"SMOKE_{stem}_extract.json").write_text(json.dumps(extracted, indent=2), encoding="utf-8")

    prs = Presentation(io.BytesIO(resp.content))
    text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ).lower()
    print(f"  slides: {len(prs.slides._sldIdLst)}")

    leaks = [
        term for term in FORBIDDEN
        if re.search(rf"(?<![a-z0-9]){re.escape(term)}(?![a-z0-9])", text)
    ]
    print(f"  leaks: {leaks or 'NONE'}")

    failed = bool(leaks)
    for kind, endpoint in [
        ("quick-wins", "/generate-quick-wins"),
        ("risk-register", "/generate-risk-register"),
        ("proposal", "/generate-proposal"),
    ]:
        start = time.time()
        pdf = client.post(f"{endpoint}?format=pdf", json=payload)
        ok = pdf.status_code == 200 and pdf.content[:5] == b"%PDF-"
        print(f"{kind:19s} {pdf.status_code}  {time.time() - start:.1f}s  "
              f"{len(pdf.content)} bytes  {'PDF ok' if ok else 'NOT A PDF'}")
        if ok:
            (out_dir / f"SMOKE_{stem}_{kind}.pdf").write_bytes(pdf.content)
        else:
            failed = True

    print("\nFAILED" if failed else "\nPASSED — artifacts in outputs/")
    return 1 if failed else 0


if __name__ == "__main__":
    notes_arg = sys.argv[1] if len(sys.argv) > 1 else "Nihaar equipments Notes.txt"
    company_arg = sys.argv[2] if len(sys.argv) > 2 else Path(notes_arg).stem.split(" Notes")[0].title()
    sys.exit(main(notes_arg, company_arg))
